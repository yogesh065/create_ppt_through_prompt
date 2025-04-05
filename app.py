import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from groq import Groq
from dotenv import load_dotenv
import os
import requests
from bs4 import BeautifulSoup
import io
import json
import base64
from PIL import Image
from io import BytesIO
import re
from urllib.parse import quote
import time
import random

# Try to import reveal_slides, with fallback if not available
try:
    import reveal_slides as rs
    REVEAL_SLIDES_AVAILABLE = True
except ImportError:
    REVEAL_SLIDES_AVAILABLE = False
    st.warning("For better slide previews, install streamlit-reveal-slides: pip install streamlit-reveal-slides")

# Set page configuration
st.set_page_config(
    page_title="Advanced Presentation Generator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better UI
st.markdown("""
<style>
    .reportview-container {
        background: #f7f7f9;
    }
    .main .block-container {
        padding-top: 2rem;
        padding-bottom: 2rem;
    }
    h1, h2, h3 {
        color: #1E3A8A;
    }
    .stButton>button {
        background-color: #1E3A8A;
        color: white;
        border-radius: 5px;
        padding: 0.5rem 1rem;
        font-weight: bold;
    }
    .stButton>button:hover {
        background-color: #2563EB;
        border-color: #2563EB;
    }
    .search-result {
        border: 1px solid #e0e0e0;
        border-radius: 5px;
        padding: 10px;
        margin-bottom: 10px;
        background: white;
    }
    .search-result h4 {
        margin-top: 0;
    }
    .theme-preview {
        border: 1px solid #ddd;
        border-radius: 5px;
        padding: 10px;
        text-align: center;
        cursor: pointer;
    }
    .theme-preview.active {
        border: 2px solid #1E3A8A;
        background-color: #f0f4ff;
    }
</style>
""", unsafe_allow_html=True)

# Load environment variables
load_dotenv()
try:
    groq_api_key = st.secrets["k"]["api_key"]
except:
    groq_api_key = os.getenv("GROQ_API_KEY")

# Initialize session state
if 'generated_content' not in st.session_state:
    st.session_state.generated_content = None
if 'slide_markdown' not in st.session_state:
    st.session_state.slide_markdown = None
if 'presentation_file' not in st.session_state:
    st.session_state.presentation_file = None
if 'search_results' not in st.session_state:
    st.session_state.search_results = None
if 'selected_theme' not in st.session_state:
    st.session_state.selected_theme = "professional"
if 'include_images' not in st.session_state:
    st.session_state.include_images = True
if 'num_slides' not in st.session_state:
    st.session_state.num_slides = 5

# Presentation themes with enhanced colors
THEMES = {
    "professional": {
        "title_font_size": Pt(36),
        "body_font_size": Pt(18),
        "title_color": RGBColor(31, 58, 138),  # Dark blue
        "accent_color": RGBColor(37, 99, 235),  # Medium blue
        "background_color": RGBColor(255, 255, 255),  # White
    },
    "minimal": {
        "title_font_size": Pt(36),
        "body_font_size": Pt(18),
        "title_color": RGBColor(30, 30, 30),  # Almost black
        "accent_color": RGBColor(100, 100, 100),  # Gray
        "background_color": RGBColor(245, 245, 245),  # Light gray
    },
    "vibrant": {
        "title_font_size": Pt(40),
        "body_font_size": Pt(20),
        "title_color": RGBColor(124, 28, 138),  # Purple
        "accent_color": RGBColor(236, 72, 153),  # Pink
        "background_color": RGBColor(253, 244, 255),  # Very light purple
    },
    "corporate": {
        "title_font_size": Pt(36),
        "body_font_size": Pt(18),
        "title_color": RGBColor(20, 83, 45),  # Dark green
        "accent_color": RGBColor(22, 163, 74),  # Green
        "background_color": RGBColor(240, 253, 244),  # Light green
    },
    "dark": {
        "title_font_size": Pt(38),
        "body_font_size": Pt(18),
        "title_color": RGBColor(226, 232, 240),  # Light gray
        "accent_color": RGBColor(56, 189, 248),  # Light blue
        "background_color": RGBColor(30, 41, 59),  # Dark blue/gray
    }
}

# Improved function to search the web with multiple fallbacks and better error handling
def search_web(query, num_results=3, max_retries=2):
    """Search the web for information related to the query with improved reliability."""
    for attempt in range(max_retries):
        try:
            # Clean and encode the query
            clean_query = quote(query)
            
            # Select a search engine - try multiple for better reliability
            search_engines = [
                # Google
                {
                    "url": f"https://www.google.com/search?q={clean_query}&num={num_results*2}",
                    "result_selector": ["div.g", "div.Gx5Zad", "div.tF2Cxc"],
                    "title_selector": ["h3", "h3.LC20lb"],
                    "link_selector": ["a"],
                    "snippet_selector": ["div.VwiC3b", "span.aCOpRe", "div.s3v9rd"]
                },
                # Bing
                {
                    "url": f"https://www.bing.com/search?q={clean_query}&count={num_results*2}",
                    "result_selector": ["li.b_algo", "div.b_title", "div.b_caption"],
                    "title_selector": ["h2", "a"],
                    "link_selector": ["a", "cite"],
                    "snippet_selector": ["p", "div.b_caption p"]
                }
            ]
            
            # Try each search engine until successful
            for engine in search_engines:
                headers = {
                    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
                }
                
                response = requests.get(engine["url"], headers=headers, timeout=10)
                
                if response.status_code != 200:
                    continue
                
                # Parse the HTML
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # Extract search results using multiple selectors for reliability
                results = []
                
                # Try each result selector
                for selector in engine["result_selector"]:
                    search_divs = soup.select(selector)
                    if search_divs:
                        break
                
                if not search_divs:
                    continue
                
                # Process found result containers
                for div in search_divs[:num_results*2]:
                    try:
                        # Try multiple title selectors
                        title = None
                        for title_selector in engine["title_selector"]:
                            title_elem = div.select_one(title_selector)
                            if title_elem:
                                title = title_elem.get_text().strip()
                                break
                        
                        if not title:
                            continue
                        
                        # Try multiple link selectors
                        link = ""
                        for link_selector in engine["link_selector"]:
                            link_elem = div.select_one(link_selector)
                            if link_elem and link_elem.has_attr('href'):
                                link = link_elem['href']
                                # Clean up Google's redirect URLs
                                if link.startswith('/url?'):
                                    link = re.search(r'url\?q=([^&]+)', link)
                                    if link:
                                        link = link.group(1)
                                break
                        
                        # Try multiple snippet selectors
                        snippet = "No description available"
                        for snippet_selector in engine["snippet_selector"]:
                            snippet_elem = div.select_one(snippet_selector)
                            if snippet_elem:
                                snippet = snippet_elem.get_text().strip()
                                break
                        
                        # Add to results if we have at least title and link
                        if title and link and link.startswith('http'):
                            results.append({
                                "title": title,
                                "link": link,
                                "snippet": snippet
                            })
                            
                            # Break once we have enough results
                            if len(results) >= num_results:
                                break
                    except Exception as e:
                        continue
                
                # If we found results, return them
                if results:
                    return results
                
                # Add delay between search engine attempts
                time.sleep(1)
            
            # If we've tried all engines and still no results, try one more approach
            try:
                # Try using DuckDuckGo as a last resort
                ddg_url = f"https://html.duckduckgo.com/html/?q={clean_query}"
                response = requests.get(ddg_url, headers=headers, timeout=10)
                
                if response.status_code == 200:
                    soup = BeautifulSoup(response.text, 'html.parser')
                    results = []
                    
                    for result in soup.select('.result'):
                        try:
                            title_elem = result.select_one('.result__title')
                            link_elem = result.select_one('.result__url')
                            snippet_elem = result.select_one('.result__snippet')
                            
                            if title_elem and link_elem:
                                title = title_elem.get_text().strip()
                                link = link_elem.get_text().strip()
                                snippet = snippet_elem.get_text().strip() if snippet_elem else "No description available"
                                
                                results.append({
                                    "title": title,
                                    "link": f"https://{link}",
                                    "snippet": snippet
                                })
                                
                                if len(results) >= num_results:
                                    break
                        except Exception:
                            continue
                    
                    if results:
                        return results
            except Exception:
                pass
            
            # Add delay between retry attempts
            if attempt < max_retries - 1:
                time.sleep(2)
        
        except Exception as e:
            if attempt < max_retries - 1:
                time.sleep(2)
    
    # If all attempts fail, return a fallback message
    return [{
        "title": "Search Failed",
        "link": "#",
        "snippet": f"Unable to retrieve search results for '{query}'. Please try again later."
    }]

# Improved function to extract content from webpages
def extract_webpage_content(url):
    try:
        if not url.startswith('http'):
            return "Invalid URL format"
            
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(url, headers=headers, timeout=10)
        
        if response.status_code != 200:
            return f"Failed to retrieve content: Status code {response.status_code}"
        
        soup = BeautifulSoup(response.text, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.extract()
        
        # Extract text from main content elements
        content_elements = soup.select('p, h1, h2, h3, h4, h5, h6, li, article, main, .content, .article')
        
        if not content_elements:
            # Fallback to extracting all text
            content = soup.get_text()
        else:
            content = ' '.join([elem.get_text() for elem in content_elements])
        
        # Clean up whitespace
        content = re.sub(r'\s+', ' ', content).strip()
        
        # Truncate if too long
        if len(content) > 5000:
            content = content[:5000] + "..."
            
        return content
    except Exception as e:
        return f"Error extracting content: {str(e)}"

# Significantly improved function to get images with multiple sources and fallbacks
def get_image_for_topic(topic, use_flowchart=False):
    """Get an image or flowchart for a given topic using multiple methods."""
    try:
        # Method 1: Use Unsplash API for reliable, high-quality images
        if not use_flowchart:
            try:
                unsplash_url = f"https://source.unsplash.com/featured/?{quote(topic)}"
                response = requests.get(unsplash_url, timeout=10)
                if response.status_code == 200:
                    # Verify it's an actual image
                    try:
                        Image.open(BytesIO(response.content))
                        return response.content
                    except:
                        pass  # Not a valid image, continue to next method
            except:
                pass  # Continue to next method if this fails
        
        # Method 2: For flowcharts, try a flowchart API/generator
        if use_flowchart:
            try:
                # Try a placeholder flowchart service
                flowchart_url = f"https://quickchart.io/graphviz?graph=digraph {{{quote(topic)}}};"
                response = requests.get(flowchart_url, timeout=10)
                if response.status_code == 200:
                    # Verify it's an actual image
                    try:
                        Image.open(BytesIO(response.content))
                        return response.content
                    except:
                        pass  # Not a valid image, continue to next method
            except:
                pass  # Continue to next method if this fails
        
        # Method 3: Try Bing image search with multiple selectors
        try:
            search_term = f"{topic} {'flowchart' if use_flowchart else ''}"
            search_url = f"https://www.bing.com/images/search?q={quote(search_term)}&first=1"
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
            }
            response = requests.get(search_url, headers=headers, timeout=10)
            
            if response.status_code == 200:
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # Try multiple selectors that might contain image URLs
                img_urls = []
                for selector in ['img.mimg', 'a.iusc img', 'img.inflnk', 'img[src^="http"]']:
                    for img in soup.select(selector):
                        if 'src' in img.attrs and img['src'].startswith('http'):
                            img_urls.append(img['src'])
                        elif 'data-src' in img.attrs and img['data-src'].startswith('http'):
                            img_urls.append(img['data-src'])
                
                # Try each URL until we find a valid image
                for img_url in img_urls[:5]:
                    try:
                        img_response = requests.get(img_url, headers=headers, timeout=5)
                        if img_response.status_code == 200:
                            # Verify it's a valid image
                            try:
                                Image.open(BytesIO(img_response.content))
                                return img_response.content
                            except:
                                continue  # Not a valid image, try next URL
                    except:
                        continue
        except:
            pass  # Continue to next method if this fails
        
        # Method 4: Generate a diagram using diagram.net API for flowcharts
        if use_flowchart:
            try:
                # Create a simple flowchart using the diagrams.net API
                flowchart_xml = f"""
                <mxGraphModel>
                    <root>
                        <mxCell id="0"/>
                        <mxCell id="1" parent="0"/>
                        <mxCell id="2" value="{topic}" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#dae8fc;strokeColor=#6c8ebf;" vertex="1" parent="1">
                            <mxGeometry x="120" y="120" width="200" height="60" as="geometry"/>
                        </mxCell>
                        <mxCell id="3" value="Process" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#d5e8d4;strokeColor=#82b366;" vertex="1" parent="1">
                            <mxGeometry x="120" y="240" width="200" height="60" as="geometry"/>
                        </mxCell>
                        <mxCell id="4" value="" style="endArrow=classic;html=1;exitX=0.5;exitY=1;exitDx=0;exitDy=0;entryX=0.5;entryY=0;entryDx=0;entryDy=0;" edge="1" parent="1" source="2" target="3">
                            <mxGeometry width="50" height="50" relative="1" as="geometry">
                                <mxPoint x="390" y="410" as="sourcePoint"/>
                                <mxPoint x="440" y="360" as="targetPoint"/>
                            </mxGeometry>
                        </mxCell>
                        <mxCell id="5" value="Output" style="rounded=1;whiteSpace=wrap;html=1;fillColor=#ffe6cc;strokeColor=#d79b00;" vertex="1" parent="1">
                            <mxGeometry x="120" y="360" width="200" height="60" as="geometry"/>
                        </mxCell>
                        <mxCell id="6" value="" style="endArrow=classic;html=1;exitX=0.5;exitY=1;exitDx=0;exitDy=0;entryX=0.5;entryY=0;entryDx=0;entryDy=0;" edge="1" parent="1" source="3" target="5">
                            <mxGeometry width="50" height="50" relative="1" as="geometry">
                                <mxPoint x="390" y="410" as="sourcePoint"/>
                                <mxPoint x="440" y="360" as="targetPoint"/>
                            </mxGeometry>
                        </mxCell>
                    </root>
                </mxGraphModel>
                """
                
                flowchart_xml_encoded = quote(flowchart_xml)
                chart_url = f"https://chart.googleapis.com/chart?cht=tx&chl={flowchart_xml_encoded}"
                response = requests.get(chart_url, timeout=10)
                
                if response.status_code == 200:
                    return response.content
            except:
                pass  # Continue to next method if this fails
        
        # Method 5: Last resort - generate a placeholder image with text
        try:
            placeholder_url = f"https://via.placeholder.com/800x600.png?text={quote(topic.replace(' ', '+'))}"
            response = requests.get(placeholder_url, timeout=10)
            if response.status_code == 200:
                return response.content
        except:
            pass
        
        # If all methods fail, return None
        return None
    except Exception as e:
        # If any unexpected error occurs, return None
        return None

# Improved function to convert presentation content to markdown for reveal.js
def pptx_to_markdown(slide_content):
    """Convert slide content to markdown for reveal.js with improved formatting."""
    markdown = "---\ntheme: black\n---\n\n"
    
    # Split slides, handling different possible delimiters
    slides = re.split(r'\n\s*\n', slide_content)
    
    for slide_text in slides:
        slide_text = slide_text.strip()
        if not slide_text:
            continue
            
        lines = slide_text.splitlines()
        if not lines:
            continue
            
        # Handle slide title - look for "Title: " prefix or just use the first line
        title_line = lines[0].strip()
        if title_line.lower().startswith("title:"):
            slide_title = title_line[6:].strip()  # Remove "Title: " prefix
        else:
            slide_title = title_line
            
        # Clean up any markdown symbols in the title
        slide_title = re.sub(r'^#+\s*', '', slide_title)  # Remove any leading # characters
        
        # Get bullet points, skipping the title line
        bullet_points = []
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
                
            # Clean up any existing bullet points to prevent doubling
            line = re.sub(r'^[-*•]\s*', '', line)
            bullet_points.append(line)
        
        # Add to markdown with proper formatting
        markdown += f"## {slide_title}\n\n"
        
        for point in bullet_points:
            markdown += f"- {point}\n"
        
        markdown += "\n---\n\n"
    
    return markdown

# Function to gather research data using multiple web searches
def gather_research_data(topic, subtopics=None):
    """Gather research data from web searches for the presentation."""
    results = {}
    
    # Main topic search
    with st.status("Searching web for information...", expanded=False) as status:
        status.update(label="Searching for main topic...")
        main_results = search_web(topic, num_results=3)
        results["main"] = main_results
        
        # Search for subtopics if provided
        if subtopics:
            subtopic_results = {}
            for i, subtopic in enumerate(subtopics):
                status.update(label=f"Researching subtopic {i+1}/{len(subtopics)}: {subtopic}")
                search_query = f"{topic} {subtopic}"
                subtopic_results[subtopic] = search_web(search_query, num_results=2)
            
            results["subtopics"] = subtopic_results
        else:
            # If no subtopics provided, generate some based on the main topic
            try:
                # Auto-generate subtopics based on main search results
                subtopics = []
                subtopic_results = {}
                
                # Extract potential subtopics from main results
                for result in main_results:
                    snippet = result.get("snippet", "")
                    title = result.get("title", "")
                    
                    # Extract phrases that might be good subtopics
                    phrases = re.findall(r'([A-Z][^.!?]*?(benefit|feature|concept|principle|type|example|use case|application)[^.!?]*)', 
                                        snippet + " " + title)
                    
                    for phrase in phrases:
                        if phrase[0] not in subtopics and len(phrase[0].split()) <= 5:
                            subtopics.append(phrase[0])
                
                # If we found some potential subtopics, search for them
                if subtopics:
                    subtopics = subtopics[:3]  # Limit to top 3 subtopics
                    for i, subtopic in enumerate(subtopics):
                        status.update(label=f"Researching auto-generated subtopic {i+1}/{len(subtopics)}: {subtopic}")
                        search_query = f"{topic} {subtopic}"
                        subtopic_results[subtopic] = search_web(search_query, num_results=2)
                    
                    results["subtopics"] = subtopic_results
            except Exception as e:
                # If auto-generation fails, just continue without subtopics
                pass
        
        # Get deeper content from the most relevant page
        if isinstance(main_results, list) and len(main_results) > 0:
            try:
                main_url = main_results[0].get("link", "")
                if main_url and main_url.startswith("http"):
                    status.update(label=f"Extracting detailed content from {main_url}")
                    results["detailed_content"] = extract_webpage_content(main_url)
            except Exception as e:
                pass
        
        status.update(label="Research completed!", state="complete")
    
    return results

# Improved function to generate slide content using Groq with research data
def groq_generate_content(topic, context, research_data, num_slides=5):
    """Generate slide content using Groq with research data."""
    if not groq_api_key:
        st.error("Please set your GROQ_API_KEY in a .env file or in Streamlit secrets.")
        return None
        
    # Initialize Groq client
    client = Groq(api_key=groq_api_key)
    
    try:
        # Format the research data for the prompt
        research_summary = "Research findings:\n"
        
        if "main" in research_data and isinstance(research_data["main"], list):
            research_summary += "Main topic search results:\n"
            for i, result in enumerate(research_data["main"][:3]):
                research_summary += f"- {result.get('title', 'No title')}: {result.get('snippet', 'No snippet')}\n"
        
        if "subtopics" in research_data:
            research_summary += "\nSubtopic search results:\n"
            for subtopic, results in research_data["subtopics"].items():
                if isinstance(results, list) and results:
                    research_summary += f"- {subtopic}: {results[0].get('snippet', 'No information')}\n"
        
        if "detailed_content" in research_data and research_data["detailed_content"]:
            content_sample = research_data["detailed_content"]
            if len(content_sample) > 1000:
                content_sample = content_sample[:1000] + "..."
            research_summary += f"\nDetailed content excerpt:\n{content_sample}\n"
        
        prompt = f"""Create a professional presentation with {num_slides} slides about "{topic}".

Use the following research data to make the presentation informative and data-driven:
{research_summary}

Additional context provided by the user: {context}

The presentation should follow these guidelines:
1. Start with a compelling title slide
2. Include an agenda or overview slide
3. Each content slide should have a clear, concise title
4. Bullet points should be specific, actionable, and data-driven using the research
5. Use the principle of "one idea per slide"
6. Include a strong concluding slide with actionable takeaways
7. Include simple flowcharts or diagrams when appropriate

IMPORTANT: For each slide, provide:
- A clear title prefixed with exactly "Title: " (this exact prefix is needed for processing)
- 3-5 concise bullet points that elaborate on the title
- Each bullet point should be on a new line without any bullet symbols (no -, *, •)
- Incorporate relevant statistics, facts, or data from the research
- Avoid jargon or overly technical terms unless necessary
- Ensure the content is engaging and visually appealing

Use this exact format for each slide:
Title: [Slide Title Here]
[Bullet point 1 - no bullet symbol]
[Bullet point 2 - no bullet symbol]
[Bullet point 3 - no bullet symbol]
[Bullet point 4 - no bullet symbol]
[Bullet point 5 - no bullet symbol]

Add a blank line between slides.

FOR FLOWCHARTS: If a slide would benefit from a simple flowchart, add a note [FLOWCHART] at the end of that slide's content.

Remember to cite sources where appropriate and maintain a professional tone."""
        
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": "You are an expert presentation designer who creates well-structured, engaging, and professional slide content backed by research data."
                },
                {
                    "role": "user",
                    "content": prompt
                }
            ],
            model="llama-3.3-70b-specdec",
            temperature=0.7,
            max_tokens=4024,
        )
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"Error generating content with Groq: {e}")
        return None

# Significantly improved function to create PowerPoint presentations with enhanced styling
def create_presentation(topic, slide_content, theme="professional", include_images=True):
    """Create a PowerPoint presentation with proper theme application and image integration."""
    prs = Presentation()
    
    # Set theme properties
    theme_properties = THEMES.get(theme, THEMES["professional"])
    
    # Function to set background color for a slide
    def apply_background(slide, color):
        """Apply background color to a slide."""
        left = top = 0
        width = prs.slide_width
        height = prs.slide_height
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.z_order = 0  # Send to back
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Apply background color to title slide
    apply_background(slide, theme_properties["background_color"])
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Apply theme to title slide
    title.text = topic
    subtitle.text = "Professional Presentation"
    
    # Apply theme formatting to title slide
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.size = theme_properties["title_font_size"]
        paragraph.font.color.rgb = theme_properties["title_color"]
        paragraph.alignment = PP_ALIGN.CENTER
    
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = theme_properties["accent_color"]
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Create content slides
    # First, clean up and parse the slide content
    slides_content = re.split(r'\n\s*\n', slide_content)
    
    for slide_index, slide_text in enumerate(slides_content):
        slide_text = slide_text.strip()
        if not slide_text:
            continue
            
        lines = slide_text.splitlines()
        if not lines:
            continue
            
        # Handle slide title
        title_line = lines[0].strip()
        if title_line.lower().startswith("title:"):
            slide_title = title_line[6:].strip()  # Remove "Title: " prefix
        else:
            slide_title = title_line
        
        # Clean up any markdown symbols in the title
        slide_title = re.sub(r'^#+\s*', '', slide_title)  # Remove any leading # characters
        
        # Get bullet points, skipping the title line
        bullet_points = []
        needs_flowchart = False
        
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
                
            # Check if the slide needs a flowchart
            if "[FLOWCHART]" in line:
                needs_flowchart = True
                line = line.replace("[FLOWCHART]", "").strip()
                
            # Clean up any existing bullet points to prevent doubling
            line = re.sub(r'^[-*•■]\s*', '', line)
            if line:
                bullet_points.append(line)
        
        # Add content slide
        content_slide_layout = prs.slide_layouts[1]  # Layout with title and content
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Apply background
        apply_background(slide, theme_properties["background_color"])
        
        # Set title
        title = slide.shapes.title
        title.text = slide_title
        
        # Apply theme formatting to title
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = theme_properties["title_font_size"]
            paragraph.font.color.rgb = theme_properties["title_color"]
            paragraph.alignment = PP_ALIGN.LEFT
        
        # Add bullet points
        if bullet_points:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = ""  # Clear any default text
            
            for point in bullet_points:
                p = tf.add_paragraph()
                p.text = point
                p.font.size = theme_properties["body_font_size"]
                p.font.color.rgb = theme_properties["accent_color"]
                p.level = 0  # First level bullet
        
        # Add an image if enabled
        if include_images:
            try:
                # Get an image related to the slide title
                image_data = get_image_for_topic(slide_title, use_flowchart=needs_flowchart)
                
                if image_data:
                    # Save the image to a BytesIO object
                    image_stream = BytesIO(image_data)
                    
                    # Add the image to the slide
                    left = Inches(7)  # Position on the right side
                    top = Inches(2)
                    width = Inches(3)  # Fixed width
                    
                    # Maintain aspect ratio
                    img = Image.open(image_stream)
                    aspect_ratio = img.height / img.width
                    height = Inches(3 * aspect_ratio)
                    
                    # Add the image to the slide
                    slide.shapes.add_picture(image_stream, left, top, width, height)
            except Exception as e:
                # Continue without an image if there was an error
                pass
    
    # Save the presentation to a BytesIO object
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io

# Main application UI with tabs
st.title("Advanced Presentation Generator")
st.markdown("### Create data-driven presentations with AI assistance and web research")

# Main content with tabs for better organization
tab1, tab2, tab3 = st.tabs(["Create Presentation", "Preview Slides", "Research Data"])

with tab1:
    col1, col2 = st.columns([2, 1])
    
    with col1:
        # Topic input
        topic = st.text_input("Presentation Topic", "", help="Enter the main topic of your presentation")
        
        # Context input
        context = st.text_area(
            "Additional Context (optional)",
            "",
            height=100,
            help="Provide any additional information or specific points to include"
        )
    
    with col2:
        # Presentation settings
        st.subheader("Presentation Settings")
        
        # Number of slides
        st.session_state.num_slides = st.slider("Number of Slides", 3, 10, 5)
        
        # Include images option
        st.session_state.include_images = st.checkbox("Include Images in Slides", value=True)
        
        # Theme selection - improved with more visual cues
        st.subheader("Select Theme")
        
        # Display theme options in a grid
        theme_cols = st.columns(3)
        
        for i, (theme_name, theme_props) in enumerate(THEMES.items()):
            with theme_cols[i % 3]:
                theme_active = st.session_state.selected_theme == theme_name
                
                # Get RGB values using tuple indexing
                background_rgb = theme_props["background_color"]
                title_rgb = theme_props["title_color"]
                
                bg_color = f"rgb({background_rgb[0]}, {background_rgb[1]}, {background_rgb[2]})"
                text_color = f"rgb({title_rgb[0]}, {title_rgb[1]}, {title_rgb[2]})"
                
                # HTML for theme preview
                st.markdown(
                    f"""
                    <div style="background-color: {bg_color}; padding: 10px; 
                        border-radius: 5px; margin-bottom: 10px;
                        border: {3 if theme_active else 1}px solid {'blue' if theme_active else '#ddd'};
                        text-align: center;">
                        <div style="color: {text_color}; font-weight: bold;">{theme_name.capitalize()}</div>
                    </div>
                    """, 
                    unsafe_allow_html=True
                )

                
                # Button to select theme
                if st.button(
                    f"Select {theme_name.capitalize()}", 
                    key=f"theme_{theme_name}",
                    type="primary" if theme_active else "secondary",
                    use_container_width=True
                ):
                    st.session_state.selected_theme = theme_name
                    st.rerun()  # Force refresh to update UI
    
    # Generate slides button
    if st.button("Generate Presentation with Web Research", use_container_width=True, type="primary"):
        if topic:
            with st.spinner("Researching and generating professional slides..."):
                # Perform web research
                research_data = gather_research_data(topic)
                st.session_state.search_results = research_data
                
                # Generate content using research data
                generated_content = groq_generate_content(
                    topic, 
                    context, 
                    research_data,
                    num_slides=st.session_state.num_slides
                )
                
                if generated_content:
                    st.session_state.generated_content = generated_content
                    
                    # Convert to markdown for preview
                    st.session_state.slide_markdown = pptx_to_markdown(generated_content)
                    
                    # Create PowerPoint file
                    pptx_io = create_presentation(
                        topic, 
                        generated_content, 
                        theme=st.session_state.selected_theme,
                        include_images=st.session_state.include_images
                    )
                    st.session_state.presentation_file = pptx_io
                    
                    st.success("Presentation generated successfully! Go to the 'Preview Slides' tab to see your presentation or check the 'Research Data' tab to view your sources.")
                else:
                    st.error("Failed to generate content. Please try again.")
        else:
            st.warning("Please enter a topic for your presentation.")
    
    # Show the generated content if available
    if st.session_state.generated_content:
        with st.expander("Generated Slide Content", expanded=True):
            edited_content = st.text_area(
                "You can edit this content before finalizing",
                st.session_state.generated_content,
                height=400
            )
            
            update_col1, update_col2 = st.columns(2)
            
            with update_col1:
                if st.button("Update Content", key="update_content", use_container_width=True):
                    # Update the stored content
                    st.session_state.generated_content = edited_content
                    
                    # Update the markdown for preview
                    st.session_state.slide_markdown = pptx_to_markdown(edited_content)
                    
                    # Update PowerPoint file
                    pptx_io = create_presentation(
                        topic, 
                        edited_content,
                        theme=st.session_state.selected_theme,
                        include_images=st.session_state.include_images
                    )
                    st.session_state.presentation_file = pptx_io
                    
                    st.success("Content updated! Go to the 'Preview Slides' tab to see your changes.")
            
            with update_col2:
                # Download button if presentation is generated
                if st.session_state.presentation_file:
                    st.download_button(
                        label="Download PowerPoint Presentation",
                        data=st.session_state.presentation_file,
                        file_name=f"{topic.replace(' ', '_')}_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                        key="download_button_tab1"
                    )

with tab2:
    st.markdown("### Preview Your Presentation")
    
    if st.session_state.slide_markdown:
        # Display presentation preview
        try:
            if REVEAL_SLIDES_AVAILABLE:
                st.markdown("#### Interactive Slide Preview")
                rs.slides(st.session_state.slide_markdown, height=500)
            else:
                # Fallback to simple preview
                st.markdown("#### Slide Content Preview")
                
                # Split into slides for better preview
                slide_parts = st.session_state.slide_markdown.split("---")
                for i, slide in enumerate(slide_parts):
                    if slide.strip():
                        with st.expander(f"Slide {i}", expanded=True):
                            st.markdown(slide)
                            
        except Exception as e:
            st.error(f"Error displaying slides: {e}")
            
            # Fallback to simple preview
            st.markdown("#### Slide Content Preview")
            st.markdown(st.session_state.slide_markdown)
            
        # Additional download button in preview tab
        if st.session_state.presentation_file:
            st.download_button(
                label="Download PowerPoint Presentation",
                data=st.session_state.presentation_file,
                file_name=f"{topic.replace(' ', '_')}_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                key="download_button_tab2"
            )
    else:
        st.info("Generate a presentation in the 'Create Presentation' tab to see a preview here.")

with tab3:
    st.markdown("### Research Data Sources")
    
    if st.session_state.search_results:
        research_data = st.session_state.search_results
        
        # Display main search results
        st.subheader("Main Topic Research")
        if "main" in research_data and isinstance(research_data["main"], list):
            for result in research_data["main"]:
                with st.container(border=True):
                    st.markdown(f"#### {result.get('title', 'No title')}")
                    st.markdown(f"**Source:** {result.get('link', 'No link')}")
                    st.markdown(f"{result.get('snippet', 'No snippet available')}")
        
        # Display subtopic results if available
        if "subtopics" in research_data:
            st.subheader("Subtopic Research")
            for subtopic, results in research_data["subtopics"].items():
                st.markdown(f"### {subtopic}")
                if isinstance(results, list):
                    for result in results:
                        with st.container(border=True):
                            st.markdown(f"#### {result.get('title', 'No title')}")
                            st.markdown(f"**Source:** {result.get('link', 'No link')}")
                            st.markdown(f"{result.get('snippet', 'No snippet available')}")
        
        # Display excerpt from detailed content if available
        if "detailed_content" in research_data and research_data["detailed_content"]:
            with st.expander("Detailed Content Excerpt"):
                st.markdown(research_data["detailed_content"][:2000] + "..." if len(research_data["detailed_content"]) > 2000 else research_data["detailed_content"])
    else:
        st.info("Generate a presentation in the 'Create Presentation' tab to see research data here.")

st.markdown("---")
st.write("Made By Yogesh Mane!")
