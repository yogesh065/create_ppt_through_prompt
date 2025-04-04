import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from groq import Groq
from dotenv import load_dotenv
import os
import requests
from bs4 import BeautifulSoup

# Load environment variables (Groq API Key)
load_dotenv()
groq_api_key = st.secrets["k"]["api_key"]
if not groq_api_key:
    st.error("Please set your GROQ_API_KEY in a .env file.")
    st.stop()

# Initialize Groq client
client = Groq(api_key=groq_api_key)

# Streamlit app title
st.title("LLM powered Slide Generator")

# Function to generate slide content using Groq
def groq_generate_content(topic, context):
    try:
        # Send request to Groq API
        response = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": f"Create 5 slide titles and bullet points about {topic}.It should be like professional slide and for each slide content should be perfectly matched.follow presentation and slide rule. Context: {context}"
                }
            ],
            model="llama-3.3-70b-specdec",  # Replace with your preferred model (e.g., llama2-70b-4096)
            temperature=0.7,
            max_tokens=4024,
        )
        generated_text = response.choices[0].message.content
        return generated_text
    except Exception as e:
        st.error(f"Error generating content with Groq: {e}")
        return None

# Function to create a PowerPoint presentation from generated content
def create_presentation(topic, slide_content):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic
    subtitle.text = "Generated using llm and Streamlit"

    # Create slides from content
    slides = slide_content.split("\n\n")  # Split into sections for each slide
    for slide_text in slides:
        lines = slide_text.split("\n")
        if len(lines) < 2:
            continue  # Skip empty or incomplete sections

        slide_title = lines[0].replace("Title: ", "")  # Extract title
        bullet_points = lines[1:]  # Extract bullet points

        # Add a new slide with bullet points
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = slide_title

        # Add bullet points to the slide
        tf = body_shape.text_frame
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point.strip()

    # Save the presentation to a file
    file_path = "generated_presentation.pptx"
    prs.save(file_path)
    return file_path

# Optional: Perform a web search for additional context using BeautifulSoup and requests
def web_search(topic):
    try:
        url = f"https://www.google.com/search?q={topic}"
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers)
        soup = BeautifulSoup(response.text, "html.parser")
        results = soup.find_all("div", class_="BNeawe vvjwJb AP7Wnd")
        return [result.text for result in results[:5]]  # Return top 5 results
    except Exception as e:
        st.error(f"Error performing web search: {e}")
        return []

# Streamlit input fields for user input
topic = st.text_input("Enter the topic of the slides:")
context = st.text_area("Enter additional context or details (optional):")

# Generate Slides Button
if st.button("Generate Slides"):
    if topic:
        with st.spinner("Generating slides using Groq..."):
            generated_content = groq_generate_content(topic, context)
            if generated_content:
                st.subheader("Generated Slide Content:")
                st.write(generated_content)

                # Create PowerPoint presentation and provide download link
                file_path = create_presentation(topic, generated_content)
                with open(file_path, "rb") as file:
                    st.download_button(
                        label="Download PowerPoint Presentation",
                        data=file,
                        file_name="generated_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
            else:
                st.error("Failed to generate content. Please try again.")
    else:
        st.warning("Please enter a topic.")

# Optional Web Search Button for Additional Context
if st.button("Perform Web Search"):
    if topic:
        with st.spinner(f"Searching the web for '{topic}'..."):
            search_results = web_search(topic)
            if search_results:
                st.subheader("Top Web Search Results:")
                for i, result in enumerate(search_results, start=1):
                    st.write(f"{i}. {result}")
            else:
                st.error("No results found.")
    else:
        st.warning("Please enter a topic to search.")

